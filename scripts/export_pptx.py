#!/usr/bin/env python3
"""Export filtered rows of data/links.csv to a PowerPoint deck.

Two modes:
  --mode table  : one slide with a table of the filtered rows
  --mode cards  : one slide per row (title, meta line, summary) for briefings

Filters: --category, --actor, --technique, --year-min
"""

import argparse
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
CSV_PATH = ROOT / "data" / "links.csv"

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x60, 0x60, 0x60)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
HEADER_BG = RGBColor(0x1A, 0x1A, 0x1A)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)


def load_rows(args):
    df = pd.read_csv(CSV_PATH, dtype=str, keep_default_na=False)

    if args.category:
        df = df[df["category"].str.strip().str.lower() == args.category.strip().lower()]
    if args.actor:
        df = df[df["actor_type"].str.strip().str.lower() == args.actor.strip().lower()]
    if args.technique:
        needle = args.technique.strip().lower()
        df = df[df["technique"].str.lower().str.contains(needle, na=False)]
    if args.year_min is not None:
        def year_ok(value):
            value = str(value).strip()
            return value.isdigit() and int(value) >= args.year_min
        df = df[df["year"].apply(year_ok)]

    def sort_key(value):
        value = str(value).strip()
        return int(value) if value.isdigit() else -1

    df = df.assign(_yk=df["year"].apply(sort_key)).sort_values("_yk", ascending=False).drop(columns="_yk")
    return df


def add_title_slide(prs, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.1), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FITD: Foot in the Door"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = INK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Physical Red Team Reference Guide"
    r2.font.size = Pt(20)
    r2.font.color.rgb = ACCENT
    if subtitle:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = subtitle
        r3.font.size = Pt(14)
        r3.font.color.rgb = MUTED


def build_table(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
    tr = title_box.text_frame.paragraphs[0].add_run()
    tr.text = "Filtered entries"
    tr.font.size = Pt(24)
    tr.font.bold = True
    tr.font.color.rgb = INK

    headers = ["Title", "Year", "Actor", "Technique", "Source"]
    rows = len(df) + 1
    table_shape = slide.shapes.add_table(
        rows, len(headers), Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )
    table = table_shape.table
    table.columns[0].width = Inches(4.3)
    table.columns[1].width = Inches(0.9)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.6)
    table.columns[4].width = Inches(2.5)

    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = head
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = HEADER_FG

    for r, (_, row) in enumerate(df.iterrows(), start=1):
        values = [
            row["title"], row["year"], row["actor_type"],
            row["technique"], row["source"],
        ]
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(value)
            run.font.size = Pt(10)
            run.font.color.rgb = INK


def build_cards(prs, df):
    for _, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.4))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tr = ttf.paragraphs[0].add_run()
        tr.text = str(row["title"])
        tr.font.size = Pt(30)
        tr.font.bold = True
        tr.font.color.rgb = INK

        meta_bits = [
            ("Actor", row["actor_type"]),
            ("Technique", row["technique"]),
            ("Year", row["year"]),
            ("Source", row["source"]),
        ]
        meta = "   |   ".join(f"{label}: {value}" for label, value in meta_bits if str(value).strip())
        meta_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.6))
        mr = meta_box.text_frame.paragraphs[0].add_run()
        mr.text = meta
        mr.font.size = Pt(14)
        mr.font.color.rgb = ACCENT

        summary_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.1), Inches(3.5))
        stf = summary_box.text_frame
        stf.word_wrap = True
        sr = stf.paragraphs[0].add_run()
        sr.text = str(row["summary"])
        sr.font.size = Pt(18)
        sr.font.color.rgb = INK

        verified = str(row["verified"]).strip().lower()
        flag_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
        fr = flag_box.text_frame.paragraphs[0].add_run()
        fr.text = f"verified: {verified}   •   {row['url']}"
        fr.font.size = Pt(11)
        fr.font.color.rgb = MUTED


def main():
    parser = argparse.ArgumentParser(description="Export filtered FITD links to a PowerPoint deck.")
    parser.add_argument("--category")
    parser.add_argument("--actor")
    parser.add_argument("--technique")
    parser.add_argument("--year-min", type=int, dest="year_min")
    parser.add_argument("--mode", choices=["table", "cards"], default="cards")
    parser.add_argument("-o", "--output", default="fitd_export.pptx")
    args = parser.parse_args()

    if not CSV_PATH.exists():
        print(f"ERROR: missing {CSV_PATH}", file=sys.stderr)
        sys.exit(1)

    df = load_rows(args)
    if df.empty:
        print("No rows match the given filters; nothing to export.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    filters = []
    for label, value in [
        ("category", args.category), ("actor", args.actor),
        ("technique", args.technique),
    ]:
        if value:
            filters.append(f"{label}={value}")
    if args.year_min is not None:
        filters.append(f"year>={args.year_min}")
    subtitle = "Filters: " + (", ".join(filters) if filters else "none") + f"  ({len(df)} rows)"

    add_title_slide(prs, subtitle)
    if args.mode == "table":
        build_table(prs, df)
    else:
        build_cards(prs, df)

    prs.save(args.output)
    print(f"OK: wrote {args.output} ({len(df)} row(s), mode={args.mode})")


if __name__ == "__main__":
    main()
